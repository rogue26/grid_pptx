from __future__ import annotations
from typing import TYPE_CHECKING, Union
from dataclasses import dataclass, field
import pandas as pd

import pptx
from pptx.chart import axis
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt

from .panel import _GridPanel, _GridPanelDefaults

# imports for type hints that would normally cause circular imports
if TYPE_CHECKING:  # pragma: no cover
    from grid_pptx import GridSlide

tick_mark_options = {
    'none': XL_TICK_MARK.NONE,
    'cross': XL_TICK_MARK.CROSS,
    'inside': XL_TICK_MARK.INSIDE,
    'outside': XL_TICK_MARK.OUTSIDE,
}

inv_tick_mark_options = {v: k for k, v in tick_mark_options.items()}

tick_label_positions = {
    'high': XL_TICK_LABEL_POSITION.HIGH,
    'low': XL_TICK_LABEL_POSITION.LOW,
    'next_to_axis': XL_TICK_LABEL_POSITION.NEXT_TO_AXIS,
    'none': XL_TICK_LABEL_POSITION.NONE,
}

inv_tick_label_positions = {v: k for k, v in tick_label_positions.items()}


@dataclass(kw_only=True)
class _ChartAxis:
    # Options for tick mark placement on either the x or y axes. Values are the constants
    # used in the python-pptx package, but are used in the code as ``XL_CHART_TYPE.LINE``, etc.
    gridchart: _GridChart
    axis_type: str

    minor_tick_marks: str = 'none'
    major_tick_marks: str = 'inside'
    has_minor_gridlines: bool = False
    has_major_gridlines: bool = False
    tick_label_position: str = 'next_to_axis'
    tick_label_italic: bool = False
    tick_label_fontsize: int = 16

    # populated later when adding the chart to the pptx slide
    axis: Union[axis.CategoryAxis, axis.DateAxis, axis.ValueAxis] = None

    def add_to_slide(self) -> None:
        # set self.axis if possible, or raise an exception if axis_type is something other than 'x' or 'y'
        if self.axis_type == 'x':
            self.axis = self.gridchart.chart.value_axis
        elif self.axis_type == 'y':
            self.axis = self.gridchart.chart.category_axis

        # set options
        self.axis.minor_tick_mark = tick_mark_options[self.minor_tick_marks]
        self.axis.major_tick_mark = tick_mark_options[self.major_tick_marks]
        self.axis.has_minor_gridlines = self.has_minor_gridlines
        self.axis.has_major_gridlines = self.has_major_gridlines
        self.axis.tick_label_position = tick_label_positions[self.tick_label_position]
        self.axis.tick_labels.font.italic = self.tick_label_italic
        self.axis.tick_labels.font.size = Pt(self.tick_label_fontsize)


class NewInitCaller(type):
    """
    metaclass that overrides the "__call__" function to automatically call "set_chart_dat_type"
    "prep_chart_data" methods after __init__, even if __init__ has been overridden
    """

    def __call__(cls, *args, **kwargs):
        """Called when you call MyNewClass() """
        obj = type.__call__(cls, *args, **kwargs)
        obj.add_axes()
        obj.set_chart_type()
        obj.prep_chart_data()
        return obj


@dataclass(kw_only=True)
class _GridChartDefaults(_GridPanelDefaults):
    chart: pptx.chart.chart.Chart = None
    x_axis: _ChartAxis = None
    y_axis: _ChartAxis = None

    # default chart parameters
    title: str = None
    has_legend: bool = True
    smooth_lines: bool = False


@dataclass(kw_only=True)
class _GridChart(_GridPanel, metaclass=NewInitCaller):
    """
    Base class for all charts

    :param df: Pandas dataframe containing data for the chart.
    :param title:
    :param has_legend:
    """
    df: pd.DataFrame
    chart_type: pptx.enum.chart.XL_CHART_TYPE = field(init=False)  # initialized in subclasses -- here as a placeholder
    chart_data: Union[CategoryChartData, XyChartData, BubbleChartData] = field(init=False)
    chart: pptx.chart.chart.Chart
    x_axis: _ChartAxis
    y_axis: _ChartAxis

    # default chart parameters
    title: str
    has_legend: bool
    smooth_lines: bool

    def add_axes(self):
        self.x_axis = _ChartAxis(gridchart=self, axis_type='x')
        self.y_axis = _ChartAxis(gridchart=self, axis_type='y')

    def evaluate_dataframe(self):
        return None

    def prep_chart_data(self) -> None:
        """

        :return:
        """
        self.chart_data.categories = self.df.index
        for column in self.df.columns:
            self.chart_data.add_series(column, self.df[column])

    def add_axes_to_slide(self):
        self.x_axis.add_to_slide()
        self.y_axis.add_to_slide()

    def add_to_slide(self, gridslide: GridSlide) -> None:
        """

        :param gridslide:
        :return:
        """

        slide = gridslide.slide
        try:
            # For some chart types, XL_CHART_TYPE.<chart type> returns a tuple with one integer. For these, we need
            # the zeroth element of the tuple

            self.chart = slide.shapes.add_chart(
                self.chart_type[0], self.x, self.y, self.cx, self.cy, self.chart_data
            ).chart

            try:
                self.chart.chart_title.text_frame.text = self.title
            except TypeError:
                pass

            self.chart.has_legend = self.has_legend
            self.chart.series[0].smooth = self.smooth_lines

            self.add_axes_to_slide()

        except TypeError:

            try:
                # For other chart types, XL_CHART_TYPE.<chart type> returns an EnumValue
                # object that is all that is needed.
                # Taking the zeroth element will throw a TypeError, which is caught here
                self.chart = slide.shapes.add_chart(
                    self.chart_type, self.x, self.y, self.cx, self.cy, self.chart_data
                ).chart

                try:
                    self.chart.chart_title.text_frame.text = self.title
                except TypeError:
                    pass

                self.chart.has_legend = self.has_legend
                self.chart.series[0].smooth = self.smooth_lines

                self.add_axes_to_slide()
            except NotImplementedError as ne:
                self.add_error_to_slide(gridslide, ne)

        except NotImplementedError as ne:
            self.add_error_to_slide(gridslide, ne)

    def add_error_to_slide(self, gridslide: GridSlide, ne: NotImplementedError) -> None:

        # Some chart types are not yet implemented in python-pptx. In those situations, attempting to add
        # the chart type to the slide will throw a NotImplementedError, which is caught here. For more info, see:
        # https://python-pptx.readthedocs.io/en/latest/dev/analysis/cht-chart-overview.html

        # for charts that are not implemented, a rectangle will be added to the slide with the error message
        # to alert the user to modify their chart choice.

        slide = gridslide.slide

        self.chart = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, self.x, self.y, self.cx, self.cy
        )

        # configure text for the rectangle
        p = self.chart.text_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        run.text = str(ne)  # text of the rectangle is the error message
        font = run.font
        font.size = Pt(10)


@dataclass(kw_only=True)
class AreaChart(_GridChartDefaults, _GridChart):
    """
    A variation of a line graph, in which areas under the line are filled in.

    :param df: Pandas dataframe containing data for the chart.
    :param three_d: Whether a 3-D version of the chart should be used
    :param stacked:
    :param normalized: Whether values should be scaled such that they sum to 100%. Only applicable \
            if ``stacked=TRUE``.

    """

    three_d: bool = False
    stacked: bool = False
    normalized: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self) -> None:

        # dictionary keys are tuples of booleans of the form (three_d, stacked, normalized)
        chart_type_dict = {
            (True, True, True): XL_CHART_TYPE.THREE_D_AREA_STACKED_100,
            (True, True, False): XL_CHART_TYPE.THREE_D_AREA_STACKED,
            # (True, False, True): 'not possible, normalized only applies to stacked',
            (True, False, False): XL_CHART_TYPE.THREE_D_AREA,
            (False, True, True): XL_CHART_TYPE.AREA_STACKED_100,
            (False, True, False): XL_CHART_TYPE.AREA_STACKED,
            # (False, False, True): 'not possible, normalized only applies to stacked',
            (False, False, False): XL_CHART_TYPE.AREA,
        }

        try:
            self.chart_type = chart_type_dict[(self.three_d, self.stacked, self.normalized)]
        except KeyError:
            raise ValueError('Charts cannot have normalized=True if stacked=False.')

        # if self.three_d:
        #     if self.stacked:
        #         if self.normalized:
        #             self.chart_type = XL_CHART_TYPE.THREE_D_AREA_STACKED_100
        #         else:
        #             self.chart_type = XL_CHART_TYPE.THREE_D_AREA_STACKED
        #     else:
        #         self.chart_type = XL_CHART_TYPE.THREE_D_AREA
        # else:
        #     if stacked:
        #         if normalized:
        #             self.chart_type = XL_CHART_TYPE.AREA_STACKED_100
        #         else:
        #             self.chart_type = XL_CHART_TYPE.AREA_STACKED
        #     else:
        #         self.chart_type = XL_CHART_TYPE.AREA


@dataclass(kw_only=True)
class BarChart(_GridChartDefaults, _GridChart):
    """
    <Description of Bar Chart>

    :param df: Pandas dataframe containing data for the chart.
    :param three_d: Whether a 3-D version of the chart should be used
    :param shape:
    :param stacked:
    :param normalized:
    """

    three_d: bool = False
    shape: str = 'rectangle'
    stacked: bool = False
    normalized: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self) -> None:
        # dictionary keys are tuples of booleans of the form (three_d, stacked, normalized)
        chart_type_dict = {
            ('rectangle', True, True, True): XL_CHART_TYPE.THREE_D_BAR_STACKED_100,
            ('rectangle', True, True, False): XL_CHART_TYPE.THREE_D_BAR_STACKED,
            # ('rectangle', True, False, True): 'not possible, normalized only applies to stacked',
            ('rectangle', True, False, False): XL_CHART_TYPE.THREE_D_BAR_CLUSTERED,
            ('rectangle', False, True, True): XL_CHART_TYPE.BAR_STACKED_100,
            ('rectangle', False, True, False): XL_CHART_TYPE.BAR_STACKED,
            # ('rectangle', False, False, True): 'not possible, normalized only applies to stacked',
            ('rectangle', False, False, False): XL_CHART_TYPE.BAR_CLUSTERED,

            ('cone', True, True, True): XL_CHART_TYPE.CONE_BAR_STACKED_100,
            ('cone', True, True, False): XL_CHART_TYPE.CONE_BAR_STACKED,
            # ('cone', True, False, True): 'not possible, normalized only applies to stacked',
            ('cone', True, False, False): XL_CHART_TYPE.CONE_BAR_CLUSTERED,
            # ('cone', False, True, True): 'not possible, cone is 3d only',
            # ('cone', False, True, False): 'not possible, cone is 3d only',
            # ('cone', False, False, True): 'not possible, normalized only applies to stacked',
            # ('cone', False, False, False): 'not possible, cone is 3d only',

            ('cylinder', True, True, True): XL_CHART_TYPE.CYLINDER_BAR_STACKED_100,
            ('cylinder', True, True, False): XL_CHART_TYPE.CYLINDER_BAR_STACKED,
            # ('cylinder', True, False, True): 'not possible, normalized only applies to stacked',
            ('cylinder', True, False, False): XL_CHART_TYPE.CYLINDER_BAR_CLUSTERED,
            # ('cylinder', False, True, True): 'not possible, cone is 3d only',
            # ('cylinder', False, True, False): 'not possible, cone is 3d only',
            # ('cylinder', False, False, True): 'not possible, normalized only applies to stacked',
            # ('cylinder', False, False, False): 'not possible, cone is 3d only',

            ('pyramid', True, True, True): XL_CHART_TYPE.PYRAMID_BAR_STACKED_100,
            ('pyramid', True, True, False): XL_CHART_TYPE.PYRAMID_BAR_STACKED,
            # ('pyramid', True, False, True): 'not possible, normalized only applies to stacked',
            ('pyramid', True, False, False): XL_CHART_TYPE.PYRAMID_BAR_CLUSTERED,
            # ('pyramid', False, True, True): 'not possible, cone is 3d only',
            # ('pyramid', False, True, False): 'not possible, cone is 3d only',
            # ('pyramid', False, False, True): 'not possible, normalized only applies to stacked',
            # ('pyramid', False, False, False): 'not possible, cone is 3d only',
        }

        try:
            self.chart_type = chart_type_dict[(self.shape, self.three_d, self.stacked, self.normalized)]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')


@dataclass(kw_only=True)
class ColumnChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param chart_data:
    :param title:
    :param has_legend:
    :param three_d: Whether a 3-D version of the chart should be used
    :param shape:
    :param stacked:
    :param normalized:
    """

    three_d: bool = False
    shape: str = 'rectangle'
    stacked: bool = False
    normalized: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self):
        chart_types = {
            'THREE_D_COLUMN': XL_CHART_TYPE.THREE_D_COLUMN,  # 3D Column.
            # 'THREE_D_COLUMN_CLUSTERED': XL_CHART_TYPE.,  # 3D Clustered Column.
            # 'THREE_D_COLUMN_STACKED': XL_CHART_TYPE.,  # 3D Stacked Column.
            # 'THREE_D_COLUMN_STACKED_100': XL_CHART_TYPE.,  # 3D 100% Stacked Column.
            # 'COLUMN_CLUSTERED': XL_CHART_TYPE.,  # Clustered Column.
            # 'COLUMN_STACKED': XL_CHART_TYPE.,  # Stacked Column.
            # 'COLUMN_STACKED_100': XL_CHART_TYPE.,  # 100% Stacked Column.
            'CONE_COL': XL_CHART_TYPE.CONE_COL,  # 3D Cone Column.
            # 'CONE_COL_CLUSTERED': XL_CHART_TYPE.,  # Clustered Cone Column.
            # 'CONE_COL_STACKED': XL_CHART_TYPE.,  # Stacked Cone Column.
            # 'CONE_COL_STACKED_100': XL_CHART_TYPE.,  # 100% Stacked Cone Column.
            # 'CYLINDER_COL': XL_CHART_TYPE.CYLINDER_COL,  # 3D Cylinder Column.
            # 'CYLINDER_COL_CLUSTERED': XL_CHART_TYPE.,  # Clustered Cone Column.
            # 'CYLINDER_COL_STACKED': XL_CHART_TYPE.,  # Stacked Cone Column.
            # 'CYLINDER_COL_STACKED_100': XL_CHART_TYPE.,  # 100% Stacked Cylinder Column.
            # 'PYRAMID_COL': XL_CHART_TYPE.PYRAMID_COL,  # 3D Pyramid Column.
            # 'PYRAMID_COL_CLUSTERED': XL_CHART_TYPE.,  # Clustered Pyramid Column.
            # 'PYRAMID_COL_STACKED': XL_CHART_TYPE.,  # Stacked Pyramid Column.
            # 'PYRAMID_COL_STACKED_100': XL_CHART_TYPE.,  # 100% Stacked Pyramid Column.
        }

        # dictionary keys are tuples of booleans of the form (three_d, stacked, normalized)
        chart_type_dict = {
            ('rectangle', True, True, True): XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100,
            ('rectangle', True, True, False): XL_CHART_TYPE.THREE_D_COLUMN_STACKED,
            # ('rectangle', True, False, True): 'not possible, normalized only applies to stacked',
            ('rectangle', True, False, False): XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED,
            ('rectangle', False, True, True): XL_CHART_TYPE.COLUMN_STACKED_100,
            ('rectangle', False, True, False): XL_CHART_TYPE.COLUMN_STACKED,
            # ('rectangle', False, False, True): 'not possible, normalized only applies to stacked',
            ('rectangle', False, False, False): XL_CHART_TYPE.COLUMN_CLUSTERED,

            ('cone', True, True, True): XL_CHART_TYPE.CONE_COL_STACKED_100,
            ('cone', True, True, False): XL_CHART_TYPE.CONE_COL_STACKED,
            # ('cone', True, False, True): 'not possible, normalized only applies to stacked',
            ('cone', True, False, False): XL_CHART_TYPE.CONE_COL_CLUSTERED,
            # ('cone', False, True, True): 'not possible, cone is 3d only',
            # ('cone', False, True, False): 'not possible, cone is 3d only',
            # ('cone', False, False, True): 'not possible, normalized only applies to stacked',
            # ('cone', False, False, False): 'not possible, cone is 3d only',

            ('cylinder', True, True, True): XL_CHART_TYPE.CYLINDER_COL_STACKED_100,
            ('cylinder', True, True, False): XL_CHART_TYPE.CYLINDER_COL_STACKED,
            # ('cylinder', True, False, True): 'not possible, normalized only applies to stacked',
            ('cylinder', True, False, False): XL_CHART_TYPE.CYLINDER_COL_CLUSTERED,
            # ('cylinder', False, True, True): 'not possible, cone is 3d only',
            # ('cylinder', False, True, False): 'not possible, cone is 3d only',
            # ('cylinder', False, False, True): 'not possible, normalized only applies to stacked',
            # ('cylinder', False, False, False): 'not possible, cone is 3d only',

            ('pyramid', True, True, True): XL_CHART_TYPE.PYRAMID_COL_STACKED_100,
            ('pyramid', True, True, False): XL_CHART_TYPE.PYRAMID_COL_STACKED,
            # ('pyramid', True, False, True): 'not possible, normalized only applies to stacked',
            ('pyramid', True, False, False): XL_CHART_TYPE.PYRAMID_COL_CLUSTERED,
            # ('pyramid', False, True, True): 'not possible, cone is 3d only',
            # ('pyramid', False, True, False): 'not possible, cone is 3d only',
            # ('pyramid', False, False, True): 'not possible, normalized only applies to stacked',
            # ('pyramid', False, False, False): 'not possible, cone is 3d only',
        }

        try:
            self.chart_type = chart_type_dict[(self.shape, self.three_d, self.stacked, self.normalized)]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')


@dataclass(kw_only=True)
class LineChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param chart_data:
    :param title:
    :param has_legend:

    :param three_d: Whether a 3-D version of the chart should be used
    :param markers:
    :param stacked:
    :param normalized:
    """

    three_d: bool = False
    markers: bool = False
    stacked: bool = False
    normalized: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self):
        # dictionary keys are tuples of booleans of the form (markers, three_d, stacked, normalized)
        chart_type_dict = {
            # (True, True, True, True): '3D with markers not currently an option',
            # (True, True, True, False): '3D with markers not currently an option',
            # (True, True, False, True): '3D with markers not currently an option',
            # (True, True, False, False): '3D with markers not currently an option',
            (True, False, True, True): XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
            (True, False, True, False): XL_CHART_TYPE.LINE_MARKERS_STACKED,
            # (True, False, False, True): 'normalized not possible unless stacked is true',
            (True, False, False, False): XL_CHART_TYPE.LINE_MARKERS,
            # (False, True, True, True): '3D version of this chart not currently available,
            # (False, True, True, False): '3D version of this chart not currently available,
            # (False, True, False, True):'3D version of this chart not currently available,
            (False, True, False, False): XL_CHART_TYPE.THREE_D_LINE,
            (False, False, True, True): XL_CHART_TYPE.LINE_STACKED_100,
            (False, False, True, False): XL_CHART_TYPE.LINE_STACKED,
            # (False, False, False, True): XL_CHART_TYPE.,
            (False, False, False, False): XL_CHART_TYPE.LINE,
        }

        try:
            self.chart_type = chart_type_dict[(self.markers, self.three_d, self.stacked, self.normalized)]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')


@dataclass(kw_only=True)
class PieChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param chart_data:
    :param title:
    :param has_legend:
    :param three_d: Whether a 3-D version of the chart should be used
    :param doughnut:
    :param exploded:
    :param compound:
    :param compound_type:
    """

    three_d: bool = False
    doughnut: bool = False
    exploded: bool = False
    compound_type: str = None

    chart_data = CategoryChartData()

    def set_chart_type(self):
        # dictionary keys are tuples of booleans of the form (three_d, exploded, doughnut, compound_type)
        chart_type_dict = {
            # (True, True, True, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (True, True, True, 'pie_of_pie'): 'This version of this chart not currently available.',
            # (True, True, True, None): 'This version of this chart not currently available.',
            # (True, True, False, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (True, True, False, 'pie_of_pie'): 'This version of this chart not currently available.',
            (True, True, False, None): XL_CHART_TYPE.THREE_D_PIE_EXPLODED,
            # (True, False, True, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (True, False, True, 'pie_of_pie'): 'This version of this chart not currently available.',
            # (True, False, True, None): 'This version of this chart not currently available.',
            # (True, False, False, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (True, False, False, 'pie_of_pie'): 'This version of this chart not currently available.',
            (True, False, False, None): XL_CHART_TYPE.THREE_D_PIE,
            # (False, True, True, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (False, True, True, 'pie_of_pie'): 'This version of this chart not currently available.',
            (False, True, True, None): XL_CHART_TYPE.DOUGHNUT_EXPLODED,
            # (False, True, False, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (False, True, False, 'pie_of_pie'): 'This version of this chart not currently available.',
            (False, True, False, None): XL_CHART_TYPE.PIE_EXPLODED,
            # (False, False, True, 'bar_of_pie'): 'This version of this chart not currently available.',
            # (False, False, True, 'pie_of_pie'): 'This version of this chart not currently available.',
            (False, False, True, None): XL_CHART_TYPE.DOUGHNUT,
            (False, False, False, 'bar_of_pie'): XL_CHART_TYPE.BAR_OF_PIE,
            (False, False, False, 'pie_of_pie'): XL_CHART_TYPE.PIE_OF_PIE,
            (False, False, False, None): XL_CHART_TYPE.PIE,
        }

        try:
            self.chart_type = chart_type_dict[
                (self.three_d, self.exploded, self.doughnut, self.compound_type)
            ]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')

    def add_axes(self):
        """ Pie charts don't have axes in python-pptx, so this method must be "nullified" """
        return None

    def add_axes_to_slide(self):
        """ Pie charts don't have axes in python-pptx, so this method must be "nullified" """
        return None


@dataclass(kw_only=True)
class RadarChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param filled:
    :param markers:
    """
    filled: bool = False
    markers: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self):
        # dictionary keys are tuples of booleans of the form (filled, markers)
        chart_type_dict = {
            # (True, True): 'This version of this chart not currently available.',
            (True, False): XL_CHART_TYPE.RADAR_FILLED,
            (False, True): XL_CHART_TYPE.RADAR_MARKERS,
            (False, False): XL_CHART_TYPE.RADAR,
        }

        try:
            self.chart_type = chart_type_dict[(self.filled, self.markers)]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')


@dataclass(kw_only=True)
class ScatterChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param x_col:
    :param y_col:
    :param lines:
    :param markers:
    :param smooth_lines:
    """

    x_col: str
    y_col: str
    lines: str = None  # options are None, 'straight', 'smooth'
    markers: bool = True

    chart_data = XyChartData()

    def set_chart_type(self):
        # dictionary keys are tuples of booleans of the form (lines, markers)
        chart_type_dict = {
            (None, True): XL_CHART_TYPE.XY_SCATTER,
            # (None, False): 'This version of this chart not currently available.',
            ('straight', True): XL_CHART_TYPE.XY_SCATTER_LINES,
            ('straight', False): XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
            ('smooth', True): XL_CHART_TYPE.XY_SCATTER_SMOOTH,
            ('smooth', False): XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
        }

        try:
            self.chart_type = chart_type_dict[(self.lines, self.markers)]
        except KeyError:
            # todo - add treatment for each possible reasons why the combation of attributes is not possible
            raise ValueError('This combination of chart attributes is not possible.')

    # def prep_chart_data(self) -> None:
    #     """
    #
    #     :return:
    #     """

    #     axis_cols = [x_col, y_col]
    #
    #     # If multi-index, how many levels, xy or bubble plot must have either 1 or 2 levels to columns
    #     if self.df.columns.nlevels == 1:
    #         # only one series in dataset
    #         series = self.chart_data.add_series('Data')
    #         for index, row in self.df.iterrows():
    #             series.add_data_point(*[row[_] for _ in axis_cols])
    #
    #     elif self.df.columns.nlevels == 2:
    #         # potentially multiple series in dataset
    #
    #         # get the list of all the axis columns (including size for bubble charts). Grid_pptx will assume that
    #         # any level 1 values in the MultiIndex that have all three of these columns should be treated as
    #         # separate series.
    #         for series_candidate in self.df.columns.get_level_values(0).unique():
    #             axis_vals_in_series = self.df.loc[:, (series_candidate, slice(None))] \
    #                 .columns.get_level_values(1).tolist()
    #
    #             # if all vals in the axis_vals_list are included in the axis_vals_in_series, then assume this
    #             # is a series and add the series and the values
    #             if all(_ in axis_vals_in_series for _ in axis_cols):
    #                 series = self.chart_data.add_series(series_candidate)
    #                 for index, row in self.df.iterrows():
    #                     series.add_data_point(*[(series_candidate, row[_]) for _ in axis_cols])
    #
    #     else:
    #         raise ValueError('The dataframe\'s columns must have no more than 2 levels.')


@dataclass(kw_only=True)
class BubbleChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param x_col:
    :param y_col:
    :param size_col:
    :param three_d: Whether a 3-D version of the chart should be used
    """

    x_col: str
    y_col: str
    size_col: str
    three_d: bool = False

    chart_data = BubbleChartData()

    def set_chart_type(self):
        if self.three_d:
            self.chart_type = XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT
        else:
            self.chart_type = XL_CHART_TYPE.BUBBLE


@dataclass(kw_only=True)
class StockChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param incl_open:
    :param volume:
    """
    incl_open: bool = False
    volume: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self):

        if self.incl_open:  #: note: "incl_open" was used instead of "open" to avoid shadowing builtin "open"
            if self.volume:
                self.chart_type = XL_CHART_TYPE.STOCK_VOHLC
            else:
                self.chart_type = XL_CHART_TYPE.STOCK_OHLC
        else:
            if self.volume:
                self.chart_type = XL_CHART_TYPE.STOCK_VHLC
            else:
                self.chart_type = XL_CHART_TYPE.STOCK_HLC


@dataclass(kw_only=True)
class SurfaceChart(_GridChartDefaults, _GridChart):
    """

    :param df: Pandas dataframe containing data for the chart.
    :param top_view:
    :param wireframe:
    """

    top_view: bool = False
    wireframe: bool = False

    chart_data = CategoryChartData()

    def set_chart_type(self):
        if self.top_view:
            if self.wireframe:
                self.chart_type = XL_CHART_TYPE.SURFACE_TOP_VIEW_WIREFRAME
            else:
                self.chart_type = XL_CHART_TYPE.SURFACE_TOP_VIEW
        else:
            if self.wireframe:
                self.chart_type = XL_CHART_TYPE.SURFACE_WIREFRAME
            else:
                self.chart_type = XL_CHART_TYPE.SURFACE
