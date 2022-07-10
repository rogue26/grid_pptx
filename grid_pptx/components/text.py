from __future__ import annotations
from typing import TYPE_CHECKING
from dataclasses import dataclass
from pptx.util import Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

from .panel import _GridPanel
from grid_pptx import colors

# imports for type hints that would normally cause circular imports
if TYPE_CHECKING:  # pragma: no cover
    from grid_pptx.slide import GridSlide

text_alignments = {
    'center': PP_PARAGRAPH_ALIGNMENT.CENTER,
    'distribute': PP_PARAGRAPH_ALIGNMENT.DISTRIBUTE,
    'justify': PP_PARAGRAPH_ALIGNMENT.JUSTIFY,
    'justify_low': PP_PARAGRAPH_ALIGNMENT.JUSTIFY_LOW,
    'left': PP_PARAGRAPH_ALIGNMENT.LEFT,
    'right': PP_PARAGRAPH_ALIGNMENT.RIGHT,
    'thai_distribute': PP_PARAGRAPH_ALIGNMENT.THAI_DISTRIBUTE,
}


@dataclass(kw_only=True)
class Text(_GridPanel):
    text: str = None
    alignment: str = 'left'
    fill_color: str = 'white'
    outline_color: str = None
    fontcolor: str = 'black'
    bold: bool = False
    fontsize: int = 16

    # def __init__(self, text: str, alignment: str = 'left', outline_color: str = None) -> None:
    #     """
    #
    #     :param text:
    #     :param alignment:
    #     """
    #     super().__init__(outline_color=outline_color)
    #
    #     self.text = text
    #
    #     self.fill_color = 'white'
    #     self.outline_color = outline_color
    #     self.fontcolor = 'black'
    #     self.bold = False
    #     self.fontsize = 16
    #     self.alignment = alignment

    def add_to_slide(self, gridslide: GridSlide) -> None:
        """

        :param gridslide:
        :return:
        """
        slide = gridslide.slide

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, self.x, self.y, self.cx, self.cy
        )

        # configure fill color
        if self.fill_color is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors.colors[self.fill_color]

        # configure outline color
        if self.outline_color is None:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.color.rgb = colors.colors[self.outline_color]

        # configure text
        p = shape.text_frame.paragraphs[0]
        p.alignment = text_alignments[self.alignment]
        run = p.add_run()
        run.text = self.text
        font = run.font
        font.size = Pt(self.fontsize)
        font.bold = self.bold
        font.color.rgb = colors.colors[self.fontcolor]


class Bullets(Text):
    def __init__(self, **kwargs) -> None:
        """

        :param kwargs:
        """
        super().__init__(**kwargs)


class Footnotes(Text):
    def __init__(self, **kwargs) -> None:
        """

        :param kwargs:
        """
        super().__init__(**kwargs)
