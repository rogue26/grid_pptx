from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


class Panel:
    color_dict = {
        'black': RGBColor(0, 0, 0),
        'white': RGBColor(255, 255, 255),
    }

    def __init__(self, slide, left, top, width, height, left_margin=0, top_margin=0, right_margin=0, bottom_margin=0,
                 row_col=None):
        self._left = left
        self._top = top
        self._width = width
        self._height = height
        self._left_margin = left_margin
        self._top_margin = top_margin
        self._right_margin = right_margin
        self._bottom_margin = bottom_margin

        # designating a panel as a row or column will inform the process of automatically subpanels based on the slide
        # design
        self._row_col = row_col

        self._slide = slide
        self.subpanels = []

    @property
    def row_col(self):
        return self._row_col

    @row_col.setter
    def row_col(self, value):
        if value in ['row', 'col']:
            self._row_col = value
        else:
            raise ValueError('Value of row_col must be either \'row\' or \'col\'')

    @property
    def left(self):
        return self._left

    @property
    def right(self):
        return self.left + self.width

    @property
    def top(self):
        return self._top

    @property
    def bottom(self):
        return self.top + self.height

    @property
    def width(self):
        return self._width

    @property
    def height(self):
        return self._height

    @property
    def slide(self):
        return self._slide

    def add_subpanel(self, obj):
        if issubclass(type(obj), Panel):
            # make sure subpanel fits within dimensions of self
            if obj.left < self.left - 0.0001:
                raise ValueError(
                    'Subpanel left side ({}) extends to the left of the parent panel left side ({})'.format(obj.left,
                                                                                                            self.left)
                )
            elif obj.right > self.right + 0.0001:
                raise ValueError(
                    'Subpanel right side ({}) extends to the right of the parent panel right side ({})'.format(
                        obj.right, self.right)
                )
            elif obj.top < self.top - 0.0001:
                raise ValueError('Subpanel top ({}) extends above the parent panel top ({})'.format(obj.top, self.top))
            elif obj.bottom > self.bottom + 0.0001:
                raise ValueError(
                    'Subpanel bottom ({}) extends below the bottom of the parent panel ({}).'.format(obj.bottom,
                                                                                                     self.bottom)
                )
            else:
                # add the subpanel
                self.subpanels.append(obj)

    def add_chart(self, df, chart_type=XL_CHART_TYPE.LINE):
        # todo: add some checks on the data

        x = Inches(self.left)
        y = Inches(self.top)
        cx = Inches(self.width)
        cy = Inches(self.height)

        chart_data = CategoryChartData()
        chart_data.categories = df.index
        for column in df.columns:
            chart_data.add_series(column, df[column])

        chart = self.slide.slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

        # format chart
        chart.has_title = False
        chart.has_legend = False
        chart.series[0].smooth = False

        # format y-axis
        chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.has_minor_gridlines = False
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE

        # format x-axis
        chart.category_axis.has_major_gridlines = False
        chart.category_axis.has_minor_gridlines = False
        chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.tick_labels.font.italic = False
        chart.category_axis.tick_labels.font.size = Pt(16)
        chart.category_axis.tick_labels.number_format = '0%'

    def add_table(self, df, header=True, firstcol=False, fontsize=12, minimize_height=True):
        # todo: add some checks on the data

        x = Inches(self.left)
        y = Inches(self.top)
        cx = Inches(self.width)

        if minimize_height:
            # setting an impossibly small initial height so the table will be as compact as possible
            cy = Inches(0.5)
        else:
            cy = Inches(self.height)

        header_adjustment = 0
        if header:
            header_adjustment = 1

        rows = len(df.index) + header_adjustment

        try:
            cols = len(df.columns)
        except AttributeError:
            # if passing a series, it will not have the columns attribute. Set cols to 1
            cols = 1

        shape = self.slide.slide.shapes.add_table(rows, cols, x, y, cx, cy)
        table = shape.table
        table.first_col = firstcol
        table.first_row = header

        if header:
            try:
                for i, val in enumerate(df.columns):
                    table.cell(0, i).text = df.columns[i]
            except AttributeError:
                # if passing a series, it will not have the columns attribute. cell value to series name
                table.cell(0, 0).text = df.index[0]

        try:
            for i, (index, row) in enumerate(df.iterrows()):
                for j, value in enumerate(row):
                    table.cell(i + header_adjustment, j).text = str(value)
                    try:
                        table.cell(i + header_adjustment, j).text_frame.paragraphs[0].runs[0].font.size = Pt(
                            fontsize)
                    except IndexError:
                        # if passing "" as the value for the table, it will raise an error when trying to set
                        # characteristics  of the font
                        pass
        except AttributeError:
            for j, value in enumerate(df):
                table.cell(j + header_adjustment, 0).text = str(value)
                try:
                    table.cell(j + header_adjustment, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(fontsize)
                except IndexError:
                    # if passing "" as the value for the table, it will raise an error when trying to set
                    # characteristics  of the font
                    pass

    def add_text(self, text, shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE, fontsize=16, bold=False, fontcolor='black',
                 outline_color=None, fill_color=None):

        # todo: add some checks on the data

        x = Inches(self.left)
        y = Inches(self.top)
        cx = Inches(self.width)
        cy = Inches(self.height)

        shape = self.slide.slide.shapes.add_shape(shape, x, y, cx, cy)

        if fill_color is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.color_dict[fill_color]

        if outline_color is None:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.color.rgb = self.color_dict[outline_color]

        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        run.text = text
        font = run.font
        font.size = Pt(fontsize)
        font.bold = bold
        font.color.rgb = self.color_dict[fontcolor]
        return shape
