import pytest
import pandas as pd
import numpy as np

from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from grid_pptx.components import table
from grid_pptx import GridPresentation, Row


@pytest.fixture
def main_df():
    return pd.DataFrame({'a': [1, 2, 3], 'b': [9, 8, 7], 'c': [7, 3, 3]})


class TestTable:
    list_of_attributes = [
        'df', 'minimize_height', 'header', 'first_col', 'fontsize', 'style'
    ]

    @pytest.fixture
    def mygridpresentation(self):
        return GridPresentation()

    @pytest.fixture
    def mytable(self, main_df, mygridpresentation):
        t = table.Table(df=main_df)
        design = Row(12, t)
        mygridpresentation.add_slide(layout_num=5, design=design, title='testing')
        return t

    def test_table_has_expected_attr(self, mytable):
        """ Test that instantiated Table object has all expected attributes

        :param mytable: PyTest fixture with an instantiated Table object
        """
        print('Missing attributes: ', *[_ for _ in self.list_of_attributes if _ not in mytable.__dict__])
        assert all(hasattr(mytable, attr) for attr in self.list_of_attributes)
