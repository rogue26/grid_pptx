import pytest
import pandas as pd
import numpy as np

from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from grid_pptx.components import chart
from grid_pptx import GridPresentation, Row


@pytest.fixture(params=[True, False])
def top_view(request):
    return request.param


@pytest.fixture(params=[True, False])
def wireframe(request):
    return request.param


@pytest.fixture(params=[True, False])
def incl_open(request):
    return request.param


@pytest.fixture(params=[True, False])
def volume(request):
    return request.param


@pytest.fixture(params=[None, 'straight', 'smooth'])
def lines(request):
    return request.param


@pytest.fixture(params=[True, False])
def filled(request):
    return request.param


@pytest.fixture(params=[None, 'bar_of_pie', 'pie_of_pie'])
def compound_type(request):
    return request.param


@pytest.fixture(params=[True, False])
def doughnut(request):
    return request.param


@pytest.fixture(params=[True, False])
def exploded(request):
    return request.param


@pytest.fixture(params=[True, False])
def markers(request):
    return request.param


@pytest.fixture(params=['rectangle', 'cone', 'cylinder', 'pyramid'])
def chart_shape(request):
    return request.param


@pytest.fixture(params=[True, False])
def three_d(request):
    return request.param


@pytest.fixture(params=[True, False])
def stacked(request):
    return request.param


@pytest.fixture(params=[True, False])
def normalized(request):
    return request.param


@pytest.fixture
def main_df():
    return pd.DataFrame({'a': [1, 2, 3], 'b': [9, 8, 7], 'c': [7, 3, 3]})


@pytest.fixture
def bad_dfs():
    """
    A list of dfs that should cause the validate_df method to throw a helpful error
    """
    dfs = [
        pd.DataFrame({'a': [1, 2, 'a'], 'b': [9, 8, 7], 'c': [7, 3, 3]}),
    ]

    return dfs


@pytest.fixture
def good_dfs():
    """
    A list of good, but possibly atypical dfs that GridChart methods should be able to work with
    """
    dfs = [
        pd.DataFrame({'a': [1, 2, 3], 'b': [9, 8, 7], 'c': [7, 3, 3]}),
        pd.DataFrame({'a': [1, 2, np.nan], 'b': [9, 8, 7], 'c': [7, 3, 3]})
    ]

    return dfs


class TestChart:
    list_of_attributes = [
        'df', 'chart_type', 'chart_data', 'x_axis', 'y_axis'
    ]

    chartclass = chart._GridChart

    @pytest.fixture
    def not_implemented(self):
        return [TestStockChart, TestSurfaceChart]

    @pytest.fixture
    def mygridpresentation(self):
        if not type(self) == TestChart:
            return GridPresentation()

    @pytest.fixture
    def mychart(self, main_df, mygridpresentation, not_implemented):
        # Note - this also has the effect of testing that these types of charts DO NOT raise a
        # NotImplemented error. Below, we'll only test that the chart types in the not_implemented
        # list DO raise an error

        if type(self) == TestScatterChart:
            c = self.chartclass(df=main_df, title='chart title', x_col='a', y_col='b')
            design = Row(12, c)
            mygridpresentation.add_slide(layout_num=5, design=design, title='testing')

        elif type(self) == TestBubbleChart:
            c = self.chartclass(df=main_df, title='chart title', x_col='a', y_col='b', size_col='c')
            design = Row(12, c)
            mygridpresentation.add_slide(layout_num=5, design=design, title='testing')

        elif not type(self) == TestChart and not type(self) in not_implemented:
            c = self.chartclass(df=main_df, title='chart title')
            design = Row(12, c)
            mygridpresentation.add_slide(layout_num=5, design=design, title='testing')

        else:
            c = None

        return c

    def test_not_implemented_error(self, main_df, mygridpresentation, not_implemented):
        if type(self) in not_implemented:
            # with pytest.raises(NotImplementedError):
            c = self.chartclass(df=main_df, title='chart title')
            design = Row(12, c)
            mygridpresentation.add_slide(layout_num=5, design=design, title='testing')

            assert c.chart.shape_type.__str__() == 'AUTO_SHAPE (1)'

    def test_chart_has_expected_attr(self, mychart, not_implemented):
        """ Test that instantiated AreaChart object has (at a minimum) all expected attributes

        :param mychart: PyTest fixture with an instantiated Chart object
        """

        if not type(self) is TestChart and type(self) not in not_implemented:
            print('Missing attributes: ', *[_ for _ in self.list_of_attributes if _ not in mychart.__dict__])
            assert all(hasattr(mychart, attr) for attr in self.list_of_attributes)

    # def test_minor_tickmarks(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting minor_tick_marks should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.minor_tick_marks = 'inside'
    #         assert mychart.x_axis.axis.minor_tick_mark == chart.tick_mark_options['inside']
    #
    #         # accessing the minor_tick_marks should return the corresponding value of the python-pptx chart
    #         assert mychart.x_axis.minor_tick_marks == 'inside'
    #
    # def test_major_tickmarks(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting major_tick_marks should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.major_tick_marks = 'inside'
    #         assert mychart.x_axis.axis.major_tick_mark == chart.tick_mark_options['inside']
    #
    #         # accessing the major_tick_marks should return the corresponding value of the python-pptx chart
    #         assert mychart.x_axis.major_tick_marks == 'inside'
    #
    # def test_minor_gridlines(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting minor gridlines should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.has_minor_gridlines = True
    #         assert mychart.x_axis.axis.has_minor_gridlines is True
    #
    #         # accessing the has_minor_gridelines attribute should return the corresponding value of the
    #         # python-pptx chart
    #         assert mychart.x_axis.has_minor_gridlines is True
    #
    #         # check the same things for False
    #         mychart.x_axis.has_minor_gridlines = False
    #         assert mychart.x_axis.axis.has_minor_gridlines is False
    #         assert mychart.x_axis.has_minor_gridlines is False
    #
    # def test_major_gridlines(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting major gridlines should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.has_major_gridlines = True
    #         assert mychart.x_axis.axis.has_major_gridlines is True
    #
    #         # accessing the has_major_gridelines attribute should return the corresponding value of the
    #         # python-pptx chart
    #         assert mychart.x_axis.has_major_gridlines is True
    #
    #         # check the same things for False
    #         mychart.x_axis.has_major_gridlines = False
    #         assert mychart.x_axis.axis.has_major_gridlines is False
    #         assert mychart.x_axis.has_major_gridlines is False
    #
    # def test_tick_label_position(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting minor_tick_marks should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.tick_label_position = 'high'
    #         assert mychart.x_axis.axis.tick_label_position == chart.tick_label_positions['high']
    #
    #         # accessing the minor_tick_marks should return the corresponding value of the python-pptx chart
    #         assert mychart.x_axis.tick_label_position == 'high'
    #
    # def test_tick_label_italic(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting tick_label_italic should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.tick_label_italic = True
    #         assert mychart.x_axis.axis.tick_labels.font.italic is True
    #
    #         # accessing the tick_label_italic attribute should return the corresponding value of the python-pptx chart
    #         assert mychart.x_axis.tick_label_italic is True
    #
    #         # check the same things for False
    #         mychart.x_axis.tick_label_italic = False
    #         assert mychart.x_axis.axis.tick_labels.font.italic is False
    #         assert mychart.x_axis.tick_label_italic is False
    #
    # def test_tick_label_fontsize(self, mychart, not_implemented):
    #     if not type(self) in [TestChart, TestPieChart] and type(self) not in not_implemented:
    #         # setting tick_label_italic should set the attribute for the chart on the python-pptx slide
    #         mychart.x_axis.tick_label_fontsize = 10
    #         assert mychart.x_axis.axis.tick_labels.font.size.pt == 10
    #
    #         # accessing the tick_label_italic attribute should return the corresponding value of the python-pptx chart
    #         assert mychart.x_axis.tick_label_fontsize == 10

    def test_chart_title(self, mychart, not_implemented):
        if not type(self) is TestChart and type(self) not in not_implemented:
            # setting title should set the attribute for the chart on the python-pptx slide
            mychart.title = 'chart title'
            assert mychart.chart.chart_title.text_frame.text == 'chart title'

            # accessing the title attribute should return the corresponding value of the python-pptx chart
            assert mychart.title == 'chart title'

    # check that chart object has only expected attributesthe
    # assert all(attr in test_chart.__dict__.keys() for attr in list_of_attributes)

    # def test_evaluate_dataframe(self):
    #     assert True
    #
    #
    # def test_set_chart_data_type(self):
    #     assert True
    #
    # def test_prep_chart_data(self):
    #     assert True
    #
    # def test_add_to_slide(self):
    #     assert True


class TestLineChart(TestChart):
    chartclass = chart.LineChart

    def test_set_chart_type(self, main_df, markers, three_d, stacked, normalized):

        if markers is True:
            if three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE_MARKERS_STACKED_100

            elif not three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE_MARKERS_STACKED

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE_MARKERS

        elif markers is False:
            if three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_LINE

            elif not three_d and stacked and normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE_STACKED_100

            elif not three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE_STACKED

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, markers=markers, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.LINE

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestAreaChart(TestChart):
    chartclass = chart.AreaChart

    def test_set_chart_type(self, main_df, three_d, stacked, normalized):

        if three_d and stacked and normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.THREE_D_AREA_STACKED_100

        elif three_d and stacked and not normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.THREE_D_AREA_STACKED

        elif three_d and not stacked and normalized:
            with pytest.raises(ValueError, match=r"Charts cannot have normalized=True if stacked=False."):
                self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)

        elif three_d and not stacked and not normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.THREE_D_AREA

        elif not three_d and stacked and normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.AREA_STACKED_100

        elif not three_d and stacked and not normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.AREA_STACKED

        elif not three_d and not stacked and normalized:
            with pytest.raises(ValueError, match=r"Charts cannot have normalized=True if stacked=False."):
                self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)

        elif not three_d and not stacked and not normalized:
            c = self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)
            assert c.chart_type == XL_CHART_TYPE.AREA

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestBarChart(TestChart):
    chartclass = chart.BarChart

    def test_set_chart_type(self, main_df, chart_shape, three_d, stacked, normalized):

        if chart_shape == 'rectangle':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_BAR_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_BAR_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_BAR_CLUSTERED

            elif not three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.BAR_STACKED_100

            elif not three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.BAR_STACKED

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)

            elif not three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.BAR_CLUSTERED

        elif chart_shape == 'cone':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_BAR_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_BAR_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_BAR_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        elif chart_shape == 'cylinder':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_BAR_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_BAR_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_BAR_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        elif chart_shape == 'pyramid':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_BAR_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_BAR_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_BAR_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestColumnChart(TestChart):
    chartclass = chart.ColumnChart

    def test_set_chart_type(self, main_df, chart_shape, three_d, stacked, normalized):

        if chart_shape == 'rectangle':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_COLUMN_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED

            elif not three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.COLUMN_STACKED_100

            elif not three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.COLUMN_STACKED

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, three_d=three_d, stacked=stacked, normalized=normalized)

            elif not three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

        elif chart_shape == 'cone':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_COL_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_COL_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CONE_COL_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        elif chart_shape == 'cylinder':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_COL_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_COL_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.CYLINDER_COL_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        elif chart_shape == 'pyramid':
            if three_d and stacked and normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_COL_STACKED_100

            elif three_d and stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_COL_STACKED

            elif three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif three_d and not stacked and not normalized:
                c = self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
                assert c.chart_type == XL_CHART_TYPE.PYRAMID_COL_CLUSTERED

            elif not three_d and stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

            elif not three_d and not stacked and normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)
            elif not three_d and not stacked and not normalized:
                with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                    self.chartclass(df=main_df, shape=chart_shape, three_d=three_d, stacked=stacked,
                                    normalized=normalized)

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestPieChart(TestChart):
    chartclass = chart.PieChart

    def test_set_chart_type(self, main_df, three_d, exploded, doughnut, compound_type):

        if three_d and exploded and doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and exploded and doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and exploded and doughnut and compound_type is None:
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and exploded and not doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and exploded and not doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and exploded and not doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.THREE_D_PIE_EXPLODED

        elif three_d and not exploded and doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and not exploded and doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and not exploded and doughnut and compound_type is None:
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and not exploded and not doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and not exploded and not doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif three_d and not exploded and not doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.THREE_D_PIE

        elif not three_d and exploded and doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and exploded and doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and exploded and doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.DOUGHNUT_EXPLODED

        elif not three_d and exploded and not doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and exploded and not doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and exploded and not doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.PIE_EXPLODED

        elif not three_d and not exploded and doughnut and compound_type == 'bar_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and not exploded and doughnut and compound_type == 'pie_of_pie':
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)

        elif not three_d and not exploded and doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.DOUGHNUT

        elif not three_d and not exploded and not doughnut and compound_type == 'bar_of_pie':
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.BAR_OF_PIE

        elif not three_d and not exploded and not doughnut and compound_type == 'pie_of_pie':
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.PIE_OF_PIE

        elif not three_d and not exploded and not doughnut and compound_type is None:
            c = self.chartclass(df=main_df, three_d=three_d, exploded=exploded, doughnut=doughnut,
                                compound_type=compound_type)
            assert c.chart_type == XL_CHART_TYPE.PIE

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestRadarChart(TestChart):
    chartclass = chart.RadarChart

    def test_set_chart_type(self, main_df, filled, markers):

        if filled and markers:
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, filled=filled, markers=markers)

        elif filled and not markers:
            c = self.chartclass(df=main_df, filled=filled, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.RADAR_FILLED

        elif not filled and markers:
            c = self.chartclass(df=main_df, filled=filled, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.RADAR_MARKERS

        elif not filled and not markers:
            c = self.chartclass(df=main_df, filled=filled, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.RADAR
        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestScatterChart(TestChart):
    chartclass = chart.ScatterChart
    list_of_attributes = [
        'df', 'chart_type', 'chart_data', 'x_axis', 'y_axis', 'x_col', 'y_col'
    ]

    def test_set_chart_type(self, main_df, lines, markers):

        if lines is None and markers:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.XY_SCATTER

        elif lines is None and not markers:
            with pytest.raises(ValueError, match=r"This combination of chart attributes is not possible."):
                self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)

        elif lines == 'straight' and markers:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.XY_SCATTER_LINES

        elif lines == 'straight' and not markers:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS

        elif lines == 'smooth' and markers:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.XY_SCATTER_SMOOTH

        elif lines == 'smooth' and not markers:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', lines=lines, markers=markers)
            assert c.chart_type == XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestBubbleChart(TestChart):
    chartclass = chart.BubbleChart
    list_of_attributes = [
        'df', 'chart_type', 'chart_data', 'x_axis', 'y_axis', 'x_col', 'y_col', 'size_col'
    ]

    def test_set_chart_type(self, main_df, three_d):

        if three_d is True:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', size_col='c', three_d=three_d)
            assert c.chart_type == XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT

        elif three_d is False:
            c = self.chartclass(df=main_df, x_col='a', y_col='b', size_col='c', three_d=three_d)
            assert c.chart_type == XL_CHART_TYPE.BUBBLE

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestStockChart(TestChart):
    chartclass = chart.StockChart

    def test_set_chart_type(self, main_df, incl_open, volume):

        if incl_open and volume:
            c = self.chartclass(df=main_df, incl_open=incl_open, volume=volume)
            assert c.chart_type == XL_CHART_TYPE.STOCK_VOHLC
        elif incl_open and not volume:
            c = self.chartclass(df=main_df, incl_open=incl_open, volume=volume)
            assert c.chart_type == XL_CHART_TYPE.STOCK_OHLC
        elif not incl_open and volume:
            c = self.chartclass(df=main_df, incl_open=incl_open, volume=volume)
            assert c.chart_type == XL_CHART_TYPE.STOCK_VHLC
        elif not incl_open and not volume:
            c = self.chartclass(df=main_df, incl_open=incl_open, volume=volume)
            assert c.chart_type == XL_CHART_TYPE.STOCK_HLC

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False


class TestSurfaceChart(TestChart):
    chartclass = chart.SurfaceChart

    def test_set_chart_type(self, main_df, top_view, wireframe):

        if top_view and wireframe:
            c = self.chartclass(df=main_df, top_view=top_view, wireframe=wireframe)
            assert c.chart_type == XL_CHART_TYPE.SURFACE_TOP_VIEW_WIREFRAME
        elif top_view and not wireframe:
            c = self.chartclass(df=main_df, top_view=top_view, wireframe=wireframe)
            assert c.chart_type == XL_CHART_TYPE.SURFACE_TOP_VIEW
        elif not top_view and wireframe:
            c = self.chartclass(df=main_df, top_view=top_view, wireframe=wireframe)
            assert c.chart_type == XL_CHART_TYPE.SURFACE_WIREFRAME
        elif not top_view and not wireframe:
            c = self.chartclass(df=main_df, top_view=top_view, wireframe=wireframe)
            assert c.chart_type == XL_CHART_TYPE.SURFACE

        else:  # pragma: no cover
            print('test not accounting for all scenarios.')
            assert False
