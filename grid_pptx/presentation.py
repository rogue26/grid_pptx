from __future__ import annotations
from typing import TYPE_CHECKING, Union
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from grid_pptx import GridSlide

# imports for type hints that would normally cause circular imports
if TYPE_CHECKING:
    from grid_pptx import Row


class GridPresentation:
    slide_sizes = {
        'letter': (10, 7.5),
        '4_3': (10, 7.5),
        '16_9': (10, 5.625),
        '16_10': (10, 6.25),
        'widescreen': (13.333, 7.5)
    }

    def __init__(self, template: Union[str, Path] = None, slide_size: Union[str, tuple, list] = None,
                 header_height: float = 1.0, footer_height: float = 0.75, left_margin: float = 0.5,
                 right_margin: float = 0.5) -> None:
        """ Stores and manipulates a pptx.Presentation instance

        :param template:
        :param slide_size:
        :param header_height:
        :param footer_height:
        :param left_margin:
        :param right_margin:
        """
        self.prs = Presentation(template)
        if template is None:
            if isinstance(slide_size, str):
                self.prs.slide_width, self.prs.slide_height = (Inches(_) for _ in self.slide_sizes[slide_size])
            elif isinstance(slide_size, (tuple, list)):
                self.prs.slide_width, self.prs.slide_height = slide_size

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        self.header_height = header_height
        self.footer_height = footer_height
        self.left_margin = left_margin
        self.right_margin = right_margin

        # self.slides = []

    def save(self, loc: Union[str, Path] = Path.home()):
        """

        :param loc:
        :return:
        """
        self.prs.save(loc)

    def add_slide(self, design: Row, layout_num: int = 5, title: str = None) -> GridSlide:
        """ Add a Slide to the pptx Presentation and a GridSlide to GridPresentation to manage it'

        :param self:
        :param design:
        :param layout_num:
        :param title:
        :return:
        """

        layout = self.prs.slide_layouts[layout_num]

        slide = self.prs.slides.add_slide(layout)
        gridslide = GridSlide(self, slide, design, title=title)
        # self.slides.append(gridslide)

        gridslide.build()
        return gridslide
