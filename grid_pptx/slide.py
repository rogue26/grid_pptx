from __future__ import annotations
from typing import TYPE_CHECKING

from pptx.slide import Slide

# imports for type hints that would normally cause circular imports
if TYPE_CHECKING:
    from grid_pptx import Row, GridPresentation


class GridSlide:

    def __init__(self, prs: GridPresentation, slide: Slide, design: Row, title=None) -> None:
        """ Manages a pptx.slides.Slide instance

        :param prs:
        :param slide:
        :param design:
        :param title:
        """
        self.slide = slide
        self.design = design

        self.left = 0.0
        self.top = 0.0
        self.width = prs.prs.slide_width.inches
        self.height = prs.prs.slide_height.inches

        self.header_height = prs.header_height
        self.footer_height = prs.footer_height
        self.left_margin = prs.left_margin
        self.right_margin = prs.right_margin

        if title is not None:
            self.title = title

    @property
    def title(self):
        """

        :return:
        """
        return self.slide.shapes.title.text

    @title.setter
    def title(self, value: str):
        """ set title for slide object

        :param value:
        :return:
        """
        self.slide.shapes.title.text = value

    def build(self):
        """

        :return:
        """
        self.design.left = self.left + self.left_margin
        self.design.width = self.width - self.left_margin - self.right_margin
        self.design.top = self.top + self.header_height
        self.design.height = self.height - self.header_height - self.footer_height

        self.design.build(self)
