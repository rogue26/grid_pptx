from __future__ import annotations
from typing import TYPE_CHECKING

from pptx.slide import Slide

# imports for type hints that would normally cause circular imports
if TYPE_CHECKING:
    from src.grid_pptx import Row, GridPresentation


class GridSlide:
    """
    Manages a pptx.slides.Slide instance
    """

    def __init__(self, prs: GridPresentation, slide: Slide, design: Row, title=None) -> None:
        self.slide = slide
        self.design = design

        self.left = 0.0
        self.top = 0.0
        self.width = prs.prs.slide_width.inches
        self.height = prs.prs.slide_height.inches

        self.header_height = prs.header_height
        self.footer_height = prs.footer_height
        self.left_margin = prs.left_margin
        self.right_margin = prs.right_margin

        if title is not None:
            self.title = title

    @property
    def title(self):
        return self.slide.shapes.title.text

    @title.setter
    def title(self, value: str):
        # set title for slide object
        self.slide.shapes.title.text = value

    def build(self):
        self.design.left = self.left + self.left_margin
        self.design.width = self.width - self.left_margin - self.right_margin
        self.design.top = self.top + self.header_height
        self.design.height = self.height - self.header_height - self.footer_height

        self.design.build(self)

    # def add_panel(self, obj):
    #     if issubclass(type(obj), GridPanel):
    #         self.panels.append(obj)
    #
    # def add_wireframe(self):
    #     def recursive(panel):
    #         for p in panel.subpanels:
    #             p.add_text("", outline_color='black')
    #             recursive(p)
    #
    #     recursive(self.panels[0])
    #
    # def loop_over_design(self, element, parent_panel):
    #     allocated_space = sum([_ for _ in element if isinstance(_, int)])
    #     unallocated_space = 12 - allocated_space
    #     number_unallocated = len([_ for _ in element if isinstance(_, list)])
    #     spaces = [_ if isinstance(_, int) else unallocated_space / number_unallocated for _ in element]
    #
    #     if parent_panel.row_col == 'row':
    #         lefts = [0] + list(itertools.accumulate([_ / 12 * parent_panel.width for _ in spaces[:-1]]))
    #         lefts = [parent_panel.left + _ for _ in lefts]
    #
    #         widths = [_ / 12 * parent_panel.width for _ in spaces]
    #         tops = [parent_panel.top for _ in spaces]
    #         heights = [parent_panel.height for _ in spaces]
    #
    #         for i, (subelement, l, t, w, h) in enumerate(zip(element, lefts, tops, widths, heights)):
    #             new_panel = GridPanel(self, l, t, w, h, row_col='col')
    #             parent_panel.add_subpanel(new_panel)
    #             if isinstance(subelement, list):
    #                 self.loop_over_design(subelement, new_panel)
    #
    #     elif parent_panel.row_col == 'col':
    #         lefts = [parent_panel.left for _ in spaces]
    #         widths = [parent_panel.width for _ in spaces]
    #
    #         tops = [0] + list(itertools.accumulate([_ / 12 * parent_panel.height for _ in spaces[:-1]]))
    #         tops = [parent_panel.top + _ for _ in tops]
    #
    #         heights = [_ / 12 * parent_panel.height for _ in spaces]
    #
    #         for (subelement, l, t, w, h) in zip(element, lefts, tops, widths, heights):
    #             new_panel = GridPanel(self, l, t, w, h, row_col='row')
    #
    #             parent_panel.add_subpanel(new_panel)
    #             if isinstance(subelement, list):
    #                 self.loop_over_design(subelement, new_panel)
    #
    # def implement_design(self):
    #     base_panel = GridPanel(self, 0, 1.2, 14.7, 6.5, row_col='row')
    #     self.add_panel(base_panel)
    #     self.loop_over_design(self.design, base_panel)
